from pptx import Presentation
from typing import Dict


class PresentationParser:
    """A class for parsing text from PowerPoint presentations."""

    def __init__(self, path_to_presentation) -> None:
        """Initializes a new instance of the pptx_parser class.
        Args:
            path_to_presentation (str): The file path of the PowerPoint presentation to parse.
        Raises:
            FileNotFoundError: If the provided file path is not valid.
        """
        self.presentation = None
        try:
            self.presentation = Presentation(path_to_presentation)
        except FileNotFoundError:
            raise FileNotFoundError("Presentation path not found.")

        # dict[slide number : text extracted from slide]
        self.presentation_text: Dict[int, str] = {}

    def parse(self):
        self.presentation_text = PresentationParser.extract_text(self.presentation)

    from pptx import Presentation
    from typing import Dict

    @staticmethod
    def extract_text(presentation: Presentation) -> Dict[int, str]:
        """
        Extracts the text from each shape in each slide of a PowerPoint presentation.

        Args:
            presentation (Presentation): The PowerPoint presentation object.

        Returns:
            Dict[int, str] of string of the slides by slide number.
        """
        presentation_text: Dict[int, str] = {}

        for slide_number, slide in enumerate(presentation.slides, start=1):
            text = PresentationParser.extract_text_from_slide(slide)

            if text.strip():
                presentation_text[slide_number] = text.strip()

        return presentation_text

    @staticmethod
    def extract_text_from_slide(slide) -> str:
        """
        Extracts the text from each shape in a slide.

        Args:
            slide: The slide object.

        Returns:
            str: The extracted text from the slide.
        """
        text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                text += PresentationParser.extract_text_from_shape(shape)

        return text.strip()

    @staticmethod
    def extract_text_from_shape(shape) -> str:
        """
        Extracts the text from a shape.

        Args:
            shape: The shape object.

        Returns:
            str: The extracted text from the shape.
        """
        text = ""

        for paragraph in shape.text_frame.paragraphs:
            text += PresentationParser.extract_text_from_paragraph(paragraph)

        return text

    @staticmethod
    def extract_text_from_paragraph(paragraph) -> str:
        """
        Extracts the text from a paragraph.

        Args:
            paragraph: The paragraph object.

        Returns:
            str: The extracted text from the paragraph.
        """
        runs_text = [run.text.strip() for run in paragraph.runs]
        text = "".join(runs_text) + "\n"

        return text

    def get_all_presentation_text(self):
        """
        get all text from presentation concatenated into one single string.
        Returns:
        str: all the text of the presentation.
        """
        return ''.join([value + ' ' for value in self.presentation_text.values()])


if __name__ == "__main__":
    pres = PresentationParser("asyncio-intro.pptx")
    pres.presentation_text = pres.extract_text(pres.presentation)
    for i in pres.presentation_text:
        print(i, ":", pres.presentation_text[i])
